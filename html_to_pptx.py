"""
將 presentation_v2.html 轉換為 PPTX 檔案
解析每一頁 slide 的文字內容，保留標題、內文、表格結構
"""

import re
import html
from pathlib import Path
from bs4 import BeautifulSoup, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 色彩定義（對應 HTML CSS 變數）──
ACCENT = RGBColor(0xD4, 0x48, 0x0B)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_BG = RGBColor(0xFF, 0xFD, 0xF7)
GRAY_BG = RGBColor(0xF5, 0xF4, 0xF0)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
ACCENT_LIGHT_BG = RGBColor(0xFF, 0xF0, 0xE8)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xEC)

# ── 投影片尺寸：16:9 ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def clean_text(el):
    """從 BeautifulSoup element 提取乾淨文字"""
    if el is None:
        return ""
    txt = el.get_text(separator=" ", strip=True)
    # 還原 HTML entities
    txt = html.unescape(txt)
    # 清理多餘空白
    txt = re.sub(r"\s+", " ", txt).strip()
    return txt


def extract_slides(html_path):
    """解析 HTML，回傳每頁投影片的結構化資料"""
    with open(html_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    slides = []
    for div in soup.find_all("div", class_="slide"):
        slide_data = {
            "type": "normal",
            "label": "",
            "title": "",
            "subtitle": "",
            "bullets": [],
            "cards": [],
            "code_blocks": [],
            "table_rows": [],
            "quote": "",
            "page_num": "",
            "is_section_divider": False,
            "is_cover": False,
        }

        classes = div.get("class", [])

        # 頁碼
        page_el = div.find("span", class_="page-num")
        if page_el:
            slide_data["page_num"] = clean_text(page_el)

        # Section divider
        if "section-divider" in classes:
            slide_data["is_section_divider"] = True
            slide_data["type"] = "section"
            label_el = div.find("div", class_="label")
            if label_el:
                slide_data["label"] = clean_text(label_el)
            h1 = div.find("h1")
            if h1:
                slide_data["title"] = clean_text(h1)
            sub_p = div.find("p", class_="subtitle")
            if sub_p:
                slide_data["subtitle"] = clean_text(sub_p)
            # 額外的 p 標籤
            for p in div.find_all("p"):
                t = clean_text(p)
                if t and t != slide_data["subtitle"] and t != slide_data["page_num"]:
                    if not slide_data["subtitle"]:
                        slide_data["subtitle"] = t
            slides.append(slide_data)
            continue

        # Cover slide
        if "slide-center" in classes:
            h1 = div.find("h1")
            if h1 and ("產險" in clean_text(h1) or "Thank" in clean_text(h1) or "謝謝" in clean_text(h1)):
                slide_data["is_cover"] = True
                slide_data["type"] = "cover"

        # Label
        label_el = div.find("div", class_="label")
        if label_el:
            slide_data["label"] = clean_text(label_el)

        # Title: h1 or h2
        h1 = div.find("h1")
        h2 = div.find("h2")
        if h1:
            slide_data["title"] = clean_text(h1)
        elif h2:
            slide_data["title"] = clean_text(h2)

        # Subtitle
        sub_p = div.find("p", class_="subtitle")
        if sub_p:
            slide_data["subtitle"] = clean_text(sub_p)

        # Cards
        for card in div.find_all("div", class_=re.compile(r"card")):
            card_classes = card.get("class", [])
            card_type = "normal"
            if "card-accent" in card_classes:
                card_type = "accent"
            elif "card-green" in card_classes:
                card_type = "green"

            card_h3 = card.find("h3")
            card_ps = card.find_all("p")
            tag_el = card.find(class_=re.compile(r"tag"))

            card_data = {
                "type": card_type,
                "tag": clean_text(tag_el) if tag_el else "",
                "title": clean_text(card_h3) if card_h3 else "",
                "text": " ".join(clean_text(p) for p in card_ps if clean_text(p)),
            }
            # 如果 card 沒有 h3，把所有文字當 text
            if not card_data["title"] and not card_data["text"]:
                card_data["text"] = clean_text(card)
            slide_data["cards"].append(card_data)

        # Bullet points (li, ol, ul)
        for li in div.find_all("li"):
            # 避免重複抓取 card 內的 li
            if li.find_parent(class_=re.compile(r"card")):
                continue
            txt = clean_text(li)
            if txt:
                slide_data["bullets"].append(txt)

        # Code blocks
        for code in div.find_all("div", class_="code-block"):
            txt = clean_text(code)
            if txt:
                slide_data["code_blocks"].append(txt)

        # Tables
        table = div.find("table")
        if table:
            rows = []
            for tr in table.find_all("tr"):
                cells = [clean_text(td) for td in tr.find_all(["th", "td"])]
                if cells:
                    rows.append(cells)
            slide_data["table_rows"] = rows

        # Quote blocks
        quote_el = div.find("div", class_="quote-block")
        if quote_el:
            slide_data["quote"] = clean_text(quote_el)

        # 如果 cards 和 bullets 都是空的，抓取 <p> 文字作為 bullets
        if not slide_data["cards"] and not slide_data["bullets"] and not slide_data["table_rows"]:
            for p in div.find_all("p"):
                txt = clean_text(p)
                if (txt and txt != slide_data["title"]
                    and txt != slide_data["subtitle"]
                    and txt != slide_data["page_num"]
                    and txt != slide_data["label"]
                    and txt != slide_data["quote"]
                    and len(txt) > 3):
                    slide_data["bullets"].append(txt)

        slides.append(slide_data)

    return slides


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Noto Sans TC"):
    """新增文字方塊"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def set_slide_bg(slide, color):
    """設定投影片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def build_cover_slide(prs, data):
    """封面頁"""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, CREAM_BG)

    # Label
    if data["label"]:
        add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
                    data["label"], font_size=12, color=ACCENT, bold=True)

    # Title
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                data["title"], font_size=44, bold=True, color=BLACK,
                alignment=PP_ALIGN.CENTER)

    # Subtitle
    if data["subtitle"]:
        add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                    data["subtitle"], font_size=22, color=GRAY,
                    alignment=PP_ALIGN.CENTER)

    # 底部資訊
    for b in data["bullets"]:
        if "侯冠宇" in b or "2025" in b:
            add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                        b, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
            break

    # 頁碼
    if data["page_num"]:
        add_textbox(slide, Inches(11), Inches(6.8), Inches(2), Inches(0.4),
                    data["page_num"], font_size=9, color=GRAY,
                    alignment=PP_ALIGN.RIGHT)


def build_section_slide(prs, data):
    """Section divider 頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, BLACK)

    if data["label"]:
        add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(0.5),
                    data["label"], font_size=12, color=ACCENT, bold=True,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.5),
                data["title"], font_size=48, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    if data["subtitle"]:
        add_textbox(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.8),
                    data["subtitle"], font_size=18, color=RGBColor(0x99, 0x99, 0x99),
                    alignment=PP_ALIGN.CENTER)

    if data["page_num"]:
        add_textbox(slide, Inches(11), Inches(6.8), Inches(2), Inches(0.4),
                    data["page_num"], font_size=9, color=RGBColor(0x66, 0x66, 0x66),
                    alignment=PP_ALIGN.RIGHT)


def build_normal_slide(prs, data):
    """一般內容頁"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, CREAM_BG)

    y_pos = Inches(0.6)
    left_margin = Inches(0.8)
    content_width = Inches(11.5)

    # Label
    if data["label"]:
        add_textbox(slide, left_margin, y_pos, Inches(3), Inches(0.35),
                    data["label"], font_size=10, color=ACCENT, bold=True)
        y_pos += Inches(0.35)

    # Title
    if data["title"]:
        add_textbox(slide, left_margin, y_pos, content_width, Inches(0.7),
                    data["title"], font_size=28, bold=True, color=BLACK)
        y_pos += Inches(0.75)

    # Subtitle
    if data["subtitle"]:
        add_textbox(slide, left_margin, y_pos, content_width, Inches(0.5),
                    data["subtitle"], font_size=16, color=GRAY)
        y_pos += Inches(0.55)

    # Table
    if data["table_rows"]:
        rows = data["table_rows"]
        n_rows = len(rows)
        n_cols = max(len(r) for r in rows) if rows else 1
        table_height = Inches(0.35 * n_rows)
        table_width = min(content_width, Inches(11))

        tbl_shape = slide.shapes.add_table(n_rows, n_cols, left_margin, y_pos,
                                           table_width, table_height)
        tbl = tbl_shape.table

        for i, row in enumerate(rows):
            for j, cell_text in enumerate(row):
                if j < n_cols:
                    cell = tbl.cell(i, j)
                    cell.text = cell_text
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(10)
                        paragraph.font.name = "Noto Sans TC"
                        if i == 0:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = WHITE
                    if i == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = BLACK

        y_pos += table_height + Inches(0.3)

    # Cards
    if data["cards"]:
        cards = data["cards"]
        # 一列最多 3 張
        per_row = min(len(cards), 3)
        card_width = Inches(11 / per_row - 0.2)
        card_height = Inches(1.2)

        for idx, card in enumerate(cards):
            col = idx % per_row
            row = idx // per_row
            x = left_margin + col * (card_width + Inches(0.2))
            y = y_pos + row * (card_height + Inches(0.15))

            # Card 背景
            bg_color = ACCENT_LIGHT_BG
            if card["type"] == "green":
                bg_color = GREEN_BG
            elif card["type"] == "normal":
                bg_color = RGBColor(0xFF, 0xFF, 0xFF)

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
            shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            shape.line.width = Pt(1)

            # Card 文字
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_top = Inches(0.1)

            card_text_parts = []
            if card["tag"]:
                card_text_parts.append(f"[{card['tag']}]")
            if card["title"]:
                card_text_parts.append(card["title"])
            if card["text"]:
                card_text_parts.append(card["text"])

            p = tf.paragraphs[0]
            full_text = "\n".join(card_text_parts)
            p.text = full_text
            p.font.size = Pt(11)
            p.font.name = "Noto Sans TC"
            p.font.color.rgb = BLACK

        total_card_rows = (len(cards) - 1) // per_row + 1
        y_pos += total_card_rows * (card_height + Inches(0.15)) + Inches(0.1)

    # Bullets
    if data["bullets"]:
        remaining = SLIDE_HEIGHT - y_pos - Inches(0.8)
        if remaining < Inches(0.5):
            remaining = Inches(2)

        txBox = slide.shapes.add_textbox(left_margin, y_pos, content_width, remaining)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(data["bullets"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(13)
            p.font.name = "Noto Sans TC"
            p.font.color.rgb = GRAY
            p.space_after = Pt(6)

    # Code blocks
    if data["code_blocks"]:
        for code in data["code_blocks"]:
            if y_pos > Inches(6.5):
                break
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left_margin, y_pos,
                content_width, Inches(0.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            shape.line.fill.background()

            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            p = tf.paragraphs[0]
            p.text = code
            p.font.size = Pt(10)
            p.font.name = "Consolas"
            p.font.color.rgb = BLACK
            y_pos += Inches(0.55)

    # Quote
    if data["quote"]:
        add_textbox(slide, left_margin, y_pos, content_width, Inches(0.5),
                    f"「{data['quote']}」", font_size=13, color=ACCENT)

    # 頁碼
    if data["page_num"]:
        add_textbox(slide, Inches(11), Inches(6.8), Inches(2), Inches(0.4),
                    data["page_num"], font_size=9, color=GRAY,
                    alignment=PP_ALIGN.RIGHT)


def main():
    base_dir = Path(__file__).parent
    html_path = base_dir / "presentation_v2.html"
    output_path = base_dir / "presentation_v2.pptx"

    print(f"讀取 HTML: {html_path}")
    slides_data = extract_slides(html_path)
    print(f"解析到 {len(slides_data)} 頁投影片")

    # 建立 PPTX
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i, data in enumerate(slides_data):
        print(f"  生成第 {i+1} 頁: [{data['type']}] {data['title'][:30]}...")
        if data["type"] == "cover" or data["is_cover"]:
            build_cover_slide(prs, data)
        elif data["type"] == "section" or data["is_section_divider"]:
            build_section_slide(prs, data)
        else:
            build_normal_slide(prs, data)

    prs.save(str(output_path))
    print(f"\n完成！已儲存至: {output_path}")
    print(f"共 {len(prs.slides)} 頁投影片")


if __name__ == "__main__":
    main()
